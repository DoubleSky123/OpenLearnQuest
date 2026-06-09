"""
Run:  pip install python-pptx
Then: python generate_poster.py
Outputs: poster.pptx  (48" x 36" landscape research poster)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colours ──────────────────────────────────────────────────────────────────
DARK_BLUE   = RGBColor(0x1a, 0x1a, 0x4e)
RED_ACCENT  = RGBColor(0xcc, 0x00, 0x00)
WHITE       = RGBColor(0xff, 0xff, 0xff)
LIGHT_BG    = RGBColor(0xf7, 0xf8, 0xfc)
BODY_TEXT   = RGBColor(0x11, 0x11, 0x11)
MUTED_TEXT  = RGBColor(0x44, 0x44, 0x44)
BORDER_COL  = RGBColor(0xdc, 0xe0, 0xef)
LIGHT_BLUE_TAG  = RGBColor(0xef, 0xf6, 0xff)
BLUE_TAG_TEXT   = RGBColor(0x1d, 0x4e, 0xd8)
RED_TAG     = RGBColor(0xfe, 0xf2, 0xf2)
RED_TAG_T   = RGBColor(0xb9, 0x1c, 0x1c)
HIGHLIGHT_BG = RGBColor(0x1a, 0x1a, 0x4e)

# ── Slide size: 48" × 36" ────────────────────────────────────────────────────
W_IN, H_IN = 48, 36
prs = Presentation()
prs.slide_width  = Inches(W_IN)
prs.slide_height = Inches(H_IN)

slide_layout = prs.slide_layouts[6]   # blank
slide = prs.slides.add_slide(slide_layout)

# ─────────────────────────────────────────────────────────────────────────────
# Helper functions
# ─────────────────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_rgb=None, line_rgb=None, line_w_pt=0):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill_rgb:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_rgb
    else:
        shape.fill.background()
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_w_pt)
    else:
        shape.line.fill.background()
    return shape


def add_text_box(slide, x, y, w, h, text, font_size, bold=False,
                  color=BODY_TEXT, align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txb


def section_box(slide, x, y, w, body_h, title_text, body_lines,
                title_h=0.35, font_body=11, bullets=False):
    """Draw a labelled section card: dark-blue title bar + light body."""
    # outer border
    add_rect(slide, x, y, w, title_h + body_h,
             fill_rgb=LIGHT_BG, line_rgb=BORDER_COL, line_w_pt=0.75)
    # title bar
    add_rect(slide, x, y, w, title_h, fill_rgb=DARK_BLUE)
    # title text
    add_text_box(slide, x+0.08, y+0.03, w-0.16, title_h-0.06,
                 title_text, 10, bold=True, color=WHITE,
                 align=PP_ALIGN.LEFT)
    # body text
    body_y = y + title_h + 0.08
    line_h = (body_h - 0.12) / max(len(body_lines), 1)
    for i, line in enumerate(body_lines):
        prefix = "• " if bullets and line.strip() else ""
        add_text_box(slide, x+0.1, body_y + i*line_h,
                     w-0.2, line_h+0.05,
                     prefix + line, font_body, color=BODY_TEXT)


# ─────────────────────────────────────────────────────────────────────────────
# HEADER  (full width, dark blue, red underline)
# ─────────────────────────────────────────────────────────────────────────────
MARGIN = 0.3
HDR_H  = 3.2

# Background
add_rect(slide, 0, 0, W_IN, HDR_H, fill_rgb=DARK_BLUE)
# Red stripe at bottom of header
add_rect(slide, 0, HDR_H - 0.18, W_IN, 0.18, fill_rgb=RED_ACCENT)

# Title
add_text_box(slide, MARGIN, 0.35, 36, 1.6,
             "OpenLearnQuest: Iterative Design of a Gamified\nAdaptive Learning Platform for Algorithm Education",
             28, bold=True, color=WHITE)

# Authors
add_text_box(slide, MARGIN, 2.05, 36, 0.5,
             "Kunyi Qiu  ·  Advisor: [Professor Name]",
             14, color=RGBColor(0xcc, 0xcc, 0xcc))
add_text_box(slide, MARGIN, 2.5, 36, 0.45,
             "Khoury College of Computer Science  ·  Northeastern University Silicon Valley",
             12, color=RGBColor(0xaa, 0xaa, 0xaa))

# NEU badge (white box, red text)
add_rect(slide, 43.5, 0.5, 4.2, 1.3, fill_rgb=WHITE,
         line_rgb=WHITE, line_w_pt=0)
add_text_box(slide, 43.5, 0.65, 4.2, 1.0,
             "NORTHEASTERN\nUNIVERSITY", 13, bold=True, color=RED_ACCENT,
             align=PP_ALIGN.CENTER)

# Poster number box
add_rect(slide, 43.8, 2.05, 3.6, 0.55, fill_rgb=RED_ACCENT)
add_text_box(slide, 43.8, 2.1, 3.6, 0.45,
             "Poster #:  ___", 12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)


# ─────────────────────────────────────────────────────────────────────────────
# LAYOUT  — three columns
# ─────────────────────────────────────────────────────────────────────────────
TOP    = HDR_H + 0.35
BOT    = H_IN - 0.3
COL_H  = BOT - TOP

COL_W  = 10.5   # left / right
MID_W  = W_IN - 2*COL_W - 2*MARGIN - 2*0.2
COL1_X = MARGIN
COL2_X = COL1_X + COL_W + 0.25
COL3_X = COL2_X + MID_W + 0.25

GAP = 0.25   # gap between sections in same column

# ─────────────────────────────────────────────────────────────────────────────
# LEFT COLUMN
# ─────────────────────────────────────────────────────────────────────────────
cy = TOP

# Background & Motivation
section_box(slide, COL1_X, cy, COL_W, 6.8,
    "BACKGROUND & MOTIVATION",
    [
        "Algorithm education presents a persistent challenge: concepts like",
        "linked lists, trees, and graphs are interconnected, yet most tools",
        "treat them in isolation. A literature review revealed a critical gap:",
        "",
        "• General learning platforms offer engagement but lack algorithm depth",
        "• Algorithm practice sites offer depth but are not novice-friendly",
        "• Educational games teach syntax, not data structure reasoning",
        "",
        "No existing platform combines gamified interaction with adaptive",
        "prerequisite routing for algorithm concept learning.",
    ], font_body=11)
cy += 6.8 + 0.35 + GAP

# Problem Statement
section_box(slide, COL1_X, cy, COL_W, 5.0,
    "PROBLEM STATEMENT",
    [
        "When students learn a complex algorithm (e.g., doubly linked list",
        "operations), they frequently discover mid-session that they lack",
        "understanding of a simpler prerequisite concept. Current platforms",
        "provide no structured mechanism to detect this gap and route students",
        "to remedial content without disrupting their learning flow.",
        "",
        "Research Question: How can research-driven iterative design produce",
        "an adaptive, gamified module that scaffolds algorithm education and",
        "enables cross-topic prerequisite routing?",
    ], font_body=11)
cy += 5.0 + 0.35 + GAP

# Literature Review
section_box(slide, COL1_X, cy, COL_W, 5.5,
    "LITERATURE REVIEW",
    [
        "Duolingo — winding path nav; XP + streak loops → mode navigation",
        "Brilliant.org — concept-first scaffolding → Tutorial intro slides",
        "Kahoot — timed challenge; competitive framing → timer & penalty",
        "CogBooks — adaptive routing on error patterns → jump logic",
        "Prodigy / Tamagotchi — companion evolution → Algo pet design",
        "BlockList (ACM SIGCSE 2025) — code-block drag-drop → core mechanic",
    ], font_body=11)
cy += 5.5 + 0.35 + GAP

# Adaptive Architecture
section_box(slide, COL1_X, cy, COL_W, 5.2,
    "ADAPTIVE LEARNING ARCHITECTURE",
    [
        "The platform models a knowledge dependency graph where modules",
        "are nodes and prerequisite relationships are directed edges:",
        "",
        "  Singly Linked List  →  Doubly Linked List  →  Trees (planned)",
        "",
        "When a student commits ≥ 2 Type-A errors (pointer logic errors",
        "shared with SLL) in a DLL session, a non-disruptive suggestion",
        "modal appears — framed as a player choice, not a system redirect.",
    ], font_body=11)


# ─────────────────────────────────────────────────────────────────────────────
# MIDDLE COLUMN
# ─────────────────────────────────────────────────────────────────────────────
cy2 = TOP

# Title bar for the big evolution section
total_mid_h = COL_H
add_rect(slide, COL2_X, cy2, MID_W, total_mid_h,
         fill_rgb=LIGHT_BG, line_rgb=BORDER_COL, line_w_pt=0.75)
add_rect(slide, COL2_X, cy2, MID_W, 0.35, fill_rgb=DARK_BLUE)
add_text_box(slide, COL2_X+0.1, cy2+0.04, MID_W-0.2, 0.28,
             "DESIGN EVOLUTION — ITERATIVE RESEARCH PROCESS",
             10, bold=True, color=WHITE)

# Six evolution steps
steps = [
    ("1", "Initial Prototype — Basic Drag-and-Drop",
     "A single-page app where students drag pseudocode blocks into an assembly area to perform singly linked list "
     "operations. Dark theme, no feedback, no distractors, no scaffolding. Confirmed the core mechanic was learnable "
     "but too abrupt for novice users.",
     "Driven by: BlockList paper + initial feasibility", "blue"),

    ("2", "Three-Mode Architecture — Tutorial → Training → Challenge",
     "Literature review showed one-size-fits-all interfaces fail novice learners. Introduced progressive scaffolding: "
     "Tutorial (step-by-step hints, fixed exercises), Training (on-demand hints, random values), and Challenge Mode "
     "(three difficulty levels, procedurally generated questions with distractor blocks and error taxonomy feedback).",
     "Driven by: Brilliant.org scaffolding model + CogBooks adaptive design", "blue"),

    ("3", "Tutorial Intro Redesign — Concept-First Onboarding",
     "Instructor feedback revealed students had no conceptual baseline on entry. Replaced the welcome popup with a "
     "4-slide concept introduction (linked list structure, node anatomy, HEAD & NULL, key terms) plus a mandatory "
     "4-question quiz. Drag-and-drop was replaced with fill-in-the-blank to reduce cognitive load.",
     "Driven by: Instructor feedback + Brilliant concept-first design", "red"),

    ("4", "Gamification Layer — Pet Companion, XP, Stars, Lives",
     "Analysis of Duolingo and Prodigy revealed long-term engagement requires emotional investment. Added a pixel "
     "husky companion (Algo) with 5 evolution stages driven by a shared XP pool, a star rating system (0–3 stars "
     "based on errors), a lives counter, and a game timer. Path-style mode navigation replaced flat card selection.",
     "Driven by: Duolingo + Tamagotchi engagement research", "blue"),

    ("5", "Doubly Linked List Module + Adaptive Cross-Module Routing",
     "Added DLL as the second module with bidirectional node visualization (⇄ arrows, prev/next states) and a new "
     "error type: Broken Prev. Cross-module adaptive logic: when ≥ 2 Type-A errors occur in a DLL session, a "
     "suggestion modal appears — framing the return to SLL as a player-initiated choice, not a system redirect.",
     "Driven by: CogBooks adaptive routing + 'player choice' design principle", "green"),

    ("6", "UI Unification, Sort Mode + Module Path Expansion",
     "Instructor usability testing identified font size and dark theme as friction points. Converted all screens to a "
     "unified light theme, doubled all font sizes, replaced drag-and-drop with click-to-place in Tutorial. Added a "
     "standalone Sort Linked List mode with live pseudocode highlighting. Module selection now shows a winding path "
     "with locked future modules (Sorting, Tree, Graph).",
     "Driven by: Instructor feedback sessions + Duolingo path navigation", "red"),
]

step_y = cy2 + 0.35 + 0.2
step_h = (total_mid_h - 0.35 - 0.3) / len(steps)

BADGE_COLORS = {"blue": RGBColor(0x1d,0x4e,0xd8), "red": RGBColor(0xb9,0x1c,0x1c), "green": RGBColor(0x15,0x80,0x3d)}
DRIVER_BG    = {"blue": LIGHT_BG,                  "red": RED_TAG,                   "green": RGBColor(0xf0,0xfd,0xf4)}

for num, title, desc, driver, color in steps:
    bx = COL2_X + 0.15
    by = step_y

    # circle badge
    badge = slide.shapes.add_shape(9, Inches(bx), Inches(by+0.05),
                                    Inches(0.38), Inches(0.38))
    badge.fill.solid()
    badge.fill.fore_color.rgb = DARK_BLUE
    badge.line.fill.background()
    tb = badge.text_frame
    tb.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tb.paragraphs[0].add_run()
    run.text = num
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = WHITE

    tx = bx + 0.5
    tw = MID_W - 0.8

    add_text_box(slide, tx, by, tw, 0.32,
                 title, 11.5, bold=True, color=DARK_BLUE)
    add_text_box(slide, tx, by+0.3, tw, step_h-0.55,
                 desc, 10.5, color=MUTED_TEXT)

    # driver tag
    drv_col = BADGE_COLORS.get(color, DARK_BLUE)
    drv_bg  = DRIVER_BG.get(color, LIGHT_BG)
    tag_y   = by + step_h - 0.32
    drv_rect = slide.shapes.add_shape(1, Inches(tx), Inches(tag_y),
                                       Inches(tw*0.75), Inches(0.26))
    drv_rect.fill.solid()
    drv_rect.fill.fore_color.rgb = drv_bg
    drv_rect.line.color.rgb = drv_col
    drv_rect.line.width = Pt(0.75)
    drv_tf = drv_rect.text_frame
    drv_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    drv_run = drv_tf.paragraphs[0].add_run()
    drv_run.text = driver
    drv_run.font.size = Pt(9)
    drv_run.font.color.rgb = drv_col

    # connector line to next
    if num != "6":
        conn = slide.shapes.add_shape(1, Inches(bx+0.17), Inches(by+step_h-0.1),
                                       Inches(0.04), Inches(0.2))
        conn.fill.solid()
        conn.fill.fore_color.rgb = BORDER_COL
        conn.line.fill.background()

    step_y += step_h


# ─────────────────────────────────────────────────────────────────────────────
# RIGHT COLUMN
# ─────────────────────────────────────────────────────────────────────────────
cy3 = TOP

# Current System Overview (with 3 stat boxes)
section_box(slide, COL3_X, cy3, COL_W, 5.5,
    "CURRENT SYSTEM OVERVIEW",
    [
        "  2 Active Modules        3 Game Modes        6 Error Types",
        "",
        "SLL Module: 4 question types × 3 levels; Tutorial, Training,",
        "            Challenge, and Sort modes",
        "",
        "DLL Module: Bidirectional visualization; Level 3 combined",
        "            operations; adaptive SLL suggestion modal",
        "",
        "Planned:    Sorting, Tree, Graph — visible as locked nodes",
        "            on the module selection path",
    ], font_body=11)
cy3 += 5.5 + 0.35 + GAP

# Error Taxonomy
section_box(slide, COL3_X, cy3, COL_W, 5.0,
    "ERROR TAXONOMY (EDUCATIONAL DESIGN)",
    [
        "Six categories map to distinct cognitive misconceptions:",
        "",
        "  Lost Reference  ·  Off-by-One  ·  NULL Pointer",
        "  Self-Loop  ·  Memory Leak  ·  Broken Prev (DLL only)",
        "",
        "Type A errors (shared with SLL) trigger adaptive routing.",
        "Type B errors (DLL-specific) trigger targeted hints only.",
        "Distractors enforce: single line, consistent variable names,",
        "and semantically meaningful wrong answers.",
    ], font_body=11)
cy3 += 5.0 + 0.35 + GAP

# Key Design Insights
section_box(slide, COL3_X, cy3, COL_W, 5.2,
    "KEY DESIGN INSIGHTS",
    [
        "Adaptive routing: Must feel like a player's decision —",
        "framing matters as much as the mechanism itself.",
        "",
        "Scaffolding order: Concept → guided practice → on-demand",
        "hints → independent challenge reduces cognitive overload.",
        "",
        "Gamification: Long-term investment (pet evolution) sustains",
        "motivation beyond immediate task rewards.",
        "",
        "Error quality: Distractors must encode real misconceptions;",
        "generic wrong answers provide no learning signal.",
    ], font_body=11)
cy3 += 5.2 + 0.35 + GAP

# Conclusions
section_box(slide, COL3_X, cy3, COL_W, 5.8,
    "CONCLUSIONS",
    [
        "This project demonstrates that a research-driven iterative",
        "process — grounded in literature review and instructor",
        "feedback — can produce an educationally principled,",
        "engaging gamified learning module.",
        "",
        "Key contributions:",
        "• Six-category error taxonomy mapping game errors to",
        "  real CS misconceptions",
        "• Cross-module adaptive routing that preserves learner",
        "  autonomy (player choice, not system redirect)",
        "• Three-tier scaffolding bridging novice-to-practitioner",
        "• Modular foundation extensible to Sorting, Tree, Graph",
    ], font_body=11)
cy3 += 5.8 + 0.35 + GAP

# References
section_box(slide, COL3_X, cy3, COL_W, 4.0,
    "REFERENCES",
    [
        "Contrino et al. (2024). CogBooks adaptive learning. Ed. Tech. Research.",
        "Duolingo Inc. (2024). Learning design principles. duolingo.com",
        "Brilliant.org (2024). Interactive learning methodology. brilliant.org",
        "Murali et al. (2025). BlockList: ACM SIGCSE TS 2025.",
        "Kim & Park (2022). Gamification in CS education. Computers & Education.",
    ], font_body=10)


# ─────────────────────────────────────────────────────────────────────────────
# Save
# ─────────────────────────────────────────────────────────────────────────────
out_path = r"D:\openLearnquest\OpenLearnQuest\poster.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
